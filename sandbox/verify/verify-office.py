#!/usr/bin/env python3
"""Exercise the shared Office CLI under the selected container runtime."""

from __future__ import annotations

import json
import subprocess
import tempfile
from pathlib import Path

from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches


def run(*args: str) -> dict:
    result = subprocess.run(
        ["artifactflow-office", *args],
        check=True,
        capture_output=True,
        text=True,
    )
    payload = json.loads(result.stdout)
    assert payload["ok"] is True, payload
    return payload


with tempfile.TemporaryDirectory(prefix="artifactflow-office-verify-") as tmp:
    root = Path(tmp)

    docx_path = root / "sample.docx"
    document = Document()
    document.add_heading("ArtifactFlow Office", level=1)
    document.add_paragraph("中文转换验证")
    document.save(docx_path)
    pdf_path = root / "sample.pdf"
    run("convert", str(docx_path), str(pdf_path))
    assert pdf_path.stat().st_size > 0

    pptx_path = root / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    text_box.text_frame.text = "Render me"
    presentation.save(pptx_path)
    pages = root / "pages"
    rendered = run("render", str(pptx_path), str(pages))
    assert rendered["pages"] == 1
    assert (pages / "page-1.png").stat().st_size > 0

    xlsx_path = root / "formula.xlsx"
    workbook = Workbook()
    worksheet = workbook.active
    worksheet["A1"] = 1
    worksheet["A2"] = 2
    worksheet["A3"] = "=SUM(A1:A2)"
    workbook.save(xlsx_path)
    recalculated_path = root / "formula-recalculated.xlsx"
    recalculated = run("recalc", str(xlsx_path), str(recalculated_path))
    assert recalculated["formulas"] == 1
    cached = load_workbook(recalculated_path, data_only=True, read_only=True)
    try:
        assert cached.active["A3"].value == 3
    finally:
        cached.close()

print("office CLI: PASS")
